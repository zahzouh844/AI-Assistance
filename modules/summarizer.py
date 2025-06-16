import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from datetime import datetime
from collections import Counter
from modules.ml_model import predict_ticket_priority
import pandas as pd
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


APP_TOKEN = "nOUdovQeRT0NTkaGlcNY6tWOPIpAf5t1L6xzJ4H0"
API_URL = "http://localhost/glpi/apirest.php"
def create_pptx_from_glpi(tickets_data, output_path, entity_name="Entité sélectionnée", users_list=None):
    if users_list is None:
        users_list = []

    prs = Presentation("template_base.pptx")
    from datetime import datetime

    # Création de la slide de titre fixe (première page)
    title_slide_layout = prs.slide_layouts[0]  # Layout "Titre"
    title_slide = prs.slides.add_slide(title_slide_layout)

    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = f" {entity_name}\n Comité technique "

    subtitle.text = " Avril/Mai 2025"


    # Déplacement manuel en première position
    prs.slides._sldIdLst.insert(0, prs.slides._sldIdLst[-1])


    # --- Slide résumé ---
    summary = summarize_tickets(tickets_data.get("tickets", []))
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Informations générales des demandes reçues"

    statuts = list(summary['statuts'].items())
    priorites = list(summary['priorites'].items())
    max_len = max(len(statuts), len(priorites), 2)

    rows = max_len + 1
    cols = 3
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 + 0.3 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # ✅ En-têtes : Infos Générales en premier
    table.cell(0, 0).text = "Informations générales"
    table.cell(0, 1).text = "Statuts"
    table.cell(0, 2).text = "Nombre de demande par priorité"

    # ✅ Colorer uniquement la première ligne (en-têtes)
    color_table_headers(table, RGBColor(95, 200, 192))
    apply_alternate_row_colors(table)

    # ✅ Colonnes Infos Générales
    table.cell(1, 0).text = f"Total tickets : {summary['nb_total']}"
    table.cell(2, 0).text = f"Dernière modification : {summary['last_mod']}"

    # ✅ Statuts et Priorités
    for i in range(max_len):
        if i < len(statuts):
            label = get_ticket_status(statuts[i][0])
            table.cell(i + 1, 1).text = f"{label} : {statuts[i][1]}"
        if i < len(priorites):
            label = get_ticket_priority(priorites[i][0])
            table.cell(i + 1, 2).text = f"{label} : {priorites[i][1]}"

    # ✅ Style général
    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(50, 50, 50)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Détails des demandes reçues "
    add_combined_ticket_table(slide, tickets_data.get("tickets", []))


    

    

    # --- Slide utilisateurs ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Comptes utilisateurs glpi actifs"

    content = f"\n Utilisateurs de l'entité ({len(users_list)}) :\n"
    for user in users_list:
        content += f"   - {user}\n"

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    tf = textbox.text_frame
    tf.text = content
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(50, 50, 50)

 
    

    # --- Slide unique : suivi mensuel global (tickets ouverts, résolus, backlog) ---
    # 1. Calcul des mois uniques
    tickets = tickets_data.get("tickets", [])
    months = sorted(set(t.get("date_creation", "")[:7] for t in tickets if t.get("date_creation")))

    # 2. Initialisation des listes
    ouverts_list = []
    resolus_list = []
    backlog_list = []

    OUVERT_STATUSES = ["Nouveau", "En cours", "En attente"]
    RESOLU_STATUSES = ["Résolu", "Clos", "Fermé"]

    cumulative_backlog = 0

    for month in months:
        ouverts = sum(1 for t in tickets if t.get('date_creation', '').startswith(month))
        resolus = sum(
            1 for t in tickets
            if get_ticket_status(t.get('status')) in RESOLU_STATUSES
            and t.get('date_mod', '').startswith(month)
        )
        cumulative_backlog += (ouverts - resolus)

        ouverts_list.append(ouverts)
        resolus_list.append(resolus)
        backlog_list.append(max(cumulative_backlog, 0))

    # --- Slide tableau : suivi mensuel global
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Suivi des incidents/demandes reçus sur la période"

    rows, cols = 4, len(months) + 1
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 + 0.3 * rows)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # En-tête
    table.cell(0, 0).text = "Suivi des incidents / demandes"
    for j, month in enumerate(months):
        mois_affiche = datetime.strptime(month, "%Y-%m").strftime("%b.-%y")
        table.cell(0, j+1).text = mois_affiche

    # Libellés lignes
    table.cell(1, 0).text = "Nombre d’incidents ouverts sur la période"
    table.cell(2, 0).text = "Nombre d’incidents résolus sur la période"
    table.cell(3, 0).text = "Nombre d’incidents en backlog (fin de mois)"

    # Remplissage des cellules
    for j in range(len(months)):
        table.cell(1, j+1).text = str(ouverts_list[j])
        table.cell(2, j+1).text = str(resolus_list[j])
        table.cell(3, j+1).text = str(backlog_list[j])

    # Mise en forme
    color_table_headers(table, RGBColor(95, 200, 192))
    apply_alternate_row_colors(table)  # <-- ici

    for row in table.rows:
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(50, 50, 50)

    # --- Slide graphique basé sur les mêmes données ---
    add_ticket_evolution_chart(prs, months, ouverts_list, resolus_list, backlog_list)

    tickets_en_cours = [t for t in tickets if t.get("status") == 2]
    # --- Autres slides ---
    sla_data = compute_sla_by_priority(tickets)
    add_sla_table_slide(prs, sla_data)
    # Ajout de la slide d’évaluation SLA
    add_sla_evaluation_table_slide(prs, tickets_en_cours)




    prs.save(output_path)









from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_ticket_evolution_chart(prs, months, ouverts_list, resolus_list, backlog_list):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Suivi des incidents / demandes "

    chart_data = CategoryChartData()
    chart_data.categories = months
    chart_data.add_series("Ouverts", ouverts_list)
    chart_data.add_series("Résolus", resolus_list)
    

    # Position et taille du graphique
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,  # ✅ Graphique à barres verticales groupées
        x, y, cx, cy, chart_data
    ).chart

    # Titre du graphique
    chart.chart_title.text_frame.text = "Évolution mensuelle des tickets"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)


    
    chart.value_axis.has_title = True
    chart.value_axis.axis_title.text_frame.text = "Nombre de tickets"

    # Légende
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT  # mettre à droite
    chart.legend.include_in_layout = False            # option recommandée
    chart.legend.font.size = Pt(12)
    

    # Couleurs personnalisées
    series_colors = {
        "Ouverts": RGBColor(0, 112, 192),    # Bleu
        "Résolus": RGBColor(0, 176, 80),     # Vert
        
    }

    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        name = series.name
        if name in series_colors:
            series.format.fill.fore_color.rgb = series_colors[name]

    # Grille légère
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(220, 220, 220)
    chart.value_axis.major_gridlines.format.line.width = Pt(0.5)







from datetime import timedelta
from dateutil import parser

from datetime import timedelta
from dateutil import parser

def normalize_priority(priority_raw):
    mapping = {
        "1": "P3",  # Très basse
        "2": "P3",  # Basse
        "3": "P2",  # Moyenne
        "4": "P1",  # Haute
        "5": "P1",  # Très haute

        "Très basse": "P3",
        "Basse": "P3",
        "Moyenne": "P2",
        "Haute": "P1",
        "Très haute": "P1",

        "P1": "P1", "P2": "P2", "P3": "P3"
    }

    return mapping.get(str(priority_raw), None)


def compute_sla_by_priority(tickets):
    sla_limits = {
        "P1": timedelta(minutes=30),
        "P2": timedelta(hours=1),
        "P3": timedelta(hours=1)
    }

    results = {
        "P1": {"total": 0, "within_sla": 0},
        "P2": {"total": 0, "within_sla": 0},
        "P3": {"total": 0, "within_sla": 0},
    }

    # Log pour vérifier le nombre de tickets extraits
    print(f"Nombre total de tickets extraits : {len(tickets)}")

    for ticket in tickets:
        priority_raw = ticket.get("priority")
        creation_str = ticket.get("date_creation")
        taken_str = ticket.get("takeintoaccountdate")

        # Vérification de la priorité
        priority = normalize_priority(priority_raw)
        if priority not in results:
            print(f"⚠️ Priorité ignorée : {priority_raw}")
            continue

        # Vérification des dates
        if not creation_str or not taken_str:
            print(f"⚠️ Dates manquantes pour ticket : {ticket}")
            continue

        try:
            creation_time = parser.parse(creation_str)
            taken_time = parser.parse(taken_str)
        except Exception as e:
            print(f"⚠️ Erreur parsing date : {e} - Ticket : {ticket}")
            continue

        # Calcul du temps de réponse
        response_time = taken_time - creation_time
        print(f"Ticket {ticket['id']} - Temps de réponse: {response_time}, Limite SLA: {sla_limits[priority]}")

        # Mise à jour des résultats
        results[priority]["total"] += 1
        if response_time <= sla_limits[priority]:
            results[priority]["within_sla"] += 1

    return results


def add_sla_table_slide(prs, sla_data):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Statistiques SLA globales"

    rows = 4
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)).table
    color_table_headers(table, RGBColor(95, 200, 192))
    apply_alternate_row_colors(table)  # <-- ici


    headers = ["Priorité", "Total Tickets", "Pris en charge dans les délais", "% Respect SLA"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h



    color_table_headers(table, RGBColor(95, 200, 192))
    apply_alternate_row_colors(table)  # <-- ici

    for idx, priority in enumerate(["P1", "P2", "P3"], start=1):
        total = sla_data[priority]["total"]
        within = sla_data[priority]["within_sla"]
        percent = f"{(within / total * 100):.1f}%" if total > 0 else "0%"

        table.cell(idx, 0).text = priority
        table.cell(idx, 1).text = str(total)
        table.cell(idx, 2).text = str(within)
        table.cell(idx, 3).text = percent




from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from datetime import datetime

def safe_parse_date(date_str):
    try:
        return datetime.strptime(date_str, "%Y-%m-%d %H:%M:%S")
    except Exception:
        return None

def calcul_sla(created, takeintoaccount, priority):
    created_dt = safe_parse_date(created)
    taken_dt = safe_parse_date(takeintoaccount)
    if not created_dt or not taken_dt:
        return "N/A"

    delay_minutes = (taken_dt - created_dt).total_seconds() / 60
    if delay_minutes < 0:
        return "N/A"

    sla_max = 30 if priority == "P1" else 60
    sla = round((sla_max / delay_minutes) * 100)
    return f"{min(sla, 100)}%"

def normalize_priority(priority_raw):
    mapping = {
        "1": "P3", "2": "P3", "3": "P2", "4": "P1", "5": "P1",
        "Très basse": "P3", "Basse": "P3", "Moyenne": "P2",
        "Haute": "P1", "Très haute": "P1", "P1": "P1", "P2": "P2", "P3": "P3"
    }

    if isinstance(priority_raw, dict):
        value = str(priority_raw.get("name") or priority_raw.get("id", "3"))
    else:
        value = str(priority_raw)

    return mapping.get(value.strip().capitalize(), "P3")

def color_table_headers(table, color):
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = color
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)
            para.alignment = PP_ALIGN.CENTER

def apply_alternate_row_colors(table):
    for i, row in enumerate(table.rows[1:], start=1):
        color = RGBColor(234, 244, 243) if i % 2 == 0 else RGBColor(245, 250, 249)
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = color

def add_sla_evaluation_table_slide(prs, tickets_en_cours):
    if not tickets_en_cours:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    slide.shapes.title.text = " Statistiques SLA par ticket"

    rows = len(tickets_en_cours) + 1
    cols = 4  # Augmentation à 4 colonnes

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.3 + rows * 0.28)


    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(1.8)

    headers = [
        "Titre du ticket",
        "Priorité",
        "Délai de prise en compte (min)",
        "SLA Respecté (%)"
    ]

    for col, header in enumerate(headers):
        table.cell(0, col).text = header

    color_table_headers(table, RGBColor(95, 200, 192))
    apply_alternate_row_colors(table)

    for i, ticket in enumerate(tickets_en_cours, start=1):
        title = ticket.get("name") or "Sans titre"
        date_creation = ticket.get("date") or ticket.get("date_creation")
        takeinto = ticket.get("takeintoaccountdate")

        created_dt = safe_parse_date(date_creation)
        taken_dt = safe_parse_date(takeinto)

        priority_raw = ticket.get("priority", "3")
        priority = normalize_priority(priority_raw)
        sla_result = calcul_sla(date_creation, takeinto, priority)

        if created_dt and taken_dt:
            delay_minutes = round((taken_dt - created_dt).total_seconds() / 60)
            delay_str = f"{delay_minutes} min"
        else:
            delay_str = "N/A"

        values = [title, priority, delay_str, sla_result]

        for j, val in enumerate(values):
            cell = table.cell(i, j)
            cell.text = val
            cell.text_frame.word_wrap = True
            cell.text_frame.margin_top = Pt(5)
            cell.text_frame.margin_bottom = Pt(5)
            cell.text_frame.margin_left = Pt(5)
            cell.text_frame.margin_right = Pt(5)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            cell.text_frame.vertical_anchor = MSO_ANCHOR.TOP

        # Ajouter les rappels SLA clairement en bas du tableau
    bottom_of_table = Inches(6.0)  # position basse fixe (ex. : 6.2 pouces depuis le haut)

    textbox_width = Inches(9)

    for idx, text in enumerate([
        " Taux de tickets (P1) pris en charge <= 30min = 100%",
        " Taux de tickets (P2) ou (P3) pris en charge <= 1h = 100%",
    ]):
        txBox = slide.shapes.add_textbox(
            left,
            bottom_of_table + Inches(0.4 * idx),  # bas de la slide
            textbox_width,
            Inches(0.4)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 70, 122)
        p.alignment = PP_ALIGN.LEFT

















from datetime import datetime, timedelta

def get_sla_alerts(tickets):
    alerts = []
    now = datetime.now()

    for ticket in tickets:
        sla_due_date = ticket.get("due_date")  # Format attendu : '2025-05-06 14:00:00'
        if not sla_due_date:
            continue

        try:
            due = datetime.strptime(sla_due_date, "%Y-%m-%d %H:%M:%S")
            delta = due - now
            if timedelta(hours=0) < delta <= timedelta(hours=1):
                alerts.append(f"⚠️ Ticket #{ticket['id']} : échéance SLA dans moins de 3h (Deadline : {due})")
            elif delta <= timedelta(seconds=0):
                alerts.append(f"❌ Ticket #{ticket['id']} : SLA DÉPASSÉE (Deadline : {due})")
        except Exception as e:
            continue

    return alerts


def format_alerts_java_style(alerts):
    if not alerts:
        return "```java\n✔️ Aucun ticket critique en attente de SLA.\n```"

    message = "```java\nALERTE SLA !\n"
    for alert in alerts:
        message += f"{alert}\n"
    message += "```"
    return message


from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LABEL_POSITION,XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches

def add_bar_chart(slide, data_dict):
    if len(data_dict) == 0:
        left, top, width, height = Inches(2), Inches(2), Inches(6), Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Aucune donnée disponible pour générer un graphique."
        return

    chart_data = CategoryChartData()
    chart_data.categories = list(data_dict.keys())
    chart_data.add_series('Tickets', list(data_dict.values()))

    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = "Répartition des tickets par type"
    chart.has_legend = False

    # Couleur claire des barres
    fill_color = RGBColor(91, 155, 213)  # bleu clair
    for series in chart.series:
        for point in series.points:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = fill_color

        # Affiche les étiquettes de données
        series.data_labels.show_value = True
        series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


import pandas as pd

def read_sdm_dashboard(filepath):
    df = pd.read_excel(filepath, sheet_name="Dashboard suivi jours 2025", header=1, engine="openpyxl")

    # Supprimer les lignes complètement vides
    df = df.dropna(how="all")

    # Supprimer les colonnes complètement vides
    df = df.dropna(axis=1, how="all")

    # Supprimer les lignes contenant "Formule à conserver"
    df = df[~df.astype(str).apply(lambda row: row.str.contains("Formule à conserver", case=False, na=False)).any(axis=1)]

    # Supprimer la colonne "Formule à conserver" si elle existe
    if "Formule à conserver" in df.columns:
        df = df.drop(columns=["Formule à conserver"])

    # Supprimer les lignes "Total" et "Reste Jetons"
    if 'Lots / Activités' in df.columns:
        df = df[~df['Lots / Activités'].astype(str).str.contains("Total|Reste Jetons", na=False)]

    # Réduire le tableau à un nombre limité de colonnes (ex : les 7 premières)
    df = df.iloc[:, :7]

    return df





    

    
def add_excel_table_slide(prs, df, title_text="Tableau de suivi OCI 2025"):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text

    rows, cols = df.shape
    max_rows = min(rows, 10)  # Ne pas surcharger le slide

    # Ajout du tableau au slide
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 + 0.3 * (max_rows + 1))  # Ajuster la hauteur selon le contenu

    table = slide.shapes.add_table(max_rows + 1, cols, left, top, width, height).table

    # Remplissage des en-têtes
    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx).text = str(col_name)

    # Remplissage des données
    for i in range(max_rows):
        for j in range(cols):
            val = df.iloc[i, j]
            table.cell(i + 1, j).text = str(val)

    # Mise en forme
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(60, 60, 60)


from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from collections import defaultdict
from modules.ml_model import predict_ticket_priority

def add_combined_ticket_table(slide, tickets):
    rows = len(tickets) + 1
    cols = 6  # 5 colonnes de base + 1 pour priorité prédite

    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(10.0)
    height = Inches(0.6 + 0.25 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table


    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # En-têtes
    headers = ["ID", "Titre", "Statut", "Priorité GLPI", "Date d'ouverture", "Priorité Prédite"]
    for j, header in enumerate(headers):
        table.cell(0, j).text = header

    # Données
    for i, ticket in enumerate(tickets):
        table.cell(i + 1, 0).text = str(ticket.get('id', 'N/A'))
        table.cell(i + 1, 1).text = ticket.get('name', 'Sans nom')
        table.cell(i + 1, 2).text = get_ticket_status(ticket.get('status'))
        table.cell(i + 1, 3).text = get_ticket_priority(ticket.get('priority'))
        table.cell(i + 1, 4).text = format_ticket_date(ticket.get('date_creation'))

        # Prédiction
        try:
            prediction = predict_ticket_priority(
                ticket.get('name'),
                ticket.get('content'),
                ticket.get('type'),
                ticket.get('itilcategories_id'),
                str(ticket.get('priority')),
                ticket.get('status'),
            )
            table.cell(i + 1, 5).text = prediction.upper()
        except Exception as e:
            table.cell(i + 1, 5).text = "Erreur"

    # Style du texte
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(50, 50, 50)

    # Couleurs
    color_table_headers(table, RGBColor(95, 200, 192))
    apply_alternate_row_colors(table)

    # ➕ Légende placée à droite du tableau
    legend_text = "P1 = Haute et très haute\nP2 = Moyenne\nP3 = Très basse et Basse"
    legend_left = left + width + Inches(0.2)  # décalé à droite du tableau
    legend_top = top + Inches(0.5)  # ajusté verticalement

    legend_shape = slide.shapes.add_textbox(legend_left, legend_top, Inches(2.5), Inches(1.0))
    text_frame = legend_shape.text_frame
    text_frame.clear()  # vide les paragraphes par défaut

    p = text_frame.paragraphs[0]
    p.text = legend_text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.alignment = PP_ALIGN.LEFT


from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from collections import defaultdict
from modules.ml_model import predict_ticket_priority

def add_predicted_priority_slide(prs, tickets, entity_name="Entité"):
    predictions = defaultdict(list)

    for ticket in tickets:
        name = ticket.get("name")
        content = ticket.get("content")
        type_ = ticket.get("type")
        category = ticket.get("itilcategories_id")
        priority = str(ticket.get("priority", ""))
        status = ticket.get("status")

        if not all([name, content, type_, status]):
            continue

        try:
            prediction = predict_ticket_priority(
                name, content, type_, category, priority, status
            )
            predictions[prediction.lower()].append(name)
        except Exception as e:
            print(f"[Erreur prédiction] Ticket '{name}' ignoré : {e}")
            continue

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"Priorités prédites par le modèle - {entity_name}"

    rows = len(predictions) + 1
    cols = 3

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # En-têtes
    table.cell(0, 0).text = "Priorité Prédite"
    table.cell(0, 1).text = "Nombre de Tickets"
    table.cell(0, 2).text = "Noms des Tickets"

    color_table_headers(table, RGBColor(95, 200, 192))
    apply_alternate_row_colors(table)

    # Ordre personnalisé : P1, P2, P3
    priority_order = ["p1", "p2", "p3"]
    sorted_predictions = sorted(
        [(p, predictions[p]) for p in predictions],
        key=lambda x: priority_order.index(x[0]) if x[0] in priority_order else 999
    )

    for i, (priority, ticket_names) in enumerate(sorted_predictions, start=1):
        table.cell(i, 0).text = priority.upper()
        table.cell(i, 1).text = str(len(ticket_names))
        table.cell(i, 2).text = "\n".join(ticket_names)

        for j in range(3):
            for paragraph in table.cell(i, j).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(40, 40, 40)  # texte en gris foncé

    # Styliser l'en-tête
    for j in range(3):
        for run in table.cell(0, j).text_frame.paragraphs[0].runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)  # texte blanc sur fond coloré








def get_glpi_tickets_by_entity_and_status(entity_id, status_id, session_token):
    url = (
        f"{API_URL}/Ticket?"
        f"criteria[0][field]=entities_id&criteria[0][searchtype]=equals&criteria[0][value]={entity_id}&"
        f"criteria[1][link]=AND&criteria[1][field]=status&criteria[1][searchtype]=equals&criteria[1][value]={status_id}&"
        f"is_recursive=0"
    )
    headers = {
        "App-Token": APP_TOKEN,
        "Session-Token": session_token
    }
    try:
        response = requests.get(url, headers=headers)
        tickets = response.json() if response.status_code == 200 else []
        # filtrage supplémentaire au besoin
        tickets = [t for t in tickets if str(t.get("entities_id")) == str(entity_id) and str(t.get("status")) == str(status_id)]
        return {"tickets": tickets}
    except Exception as e:
        return {"error": f"Erreur API : {e}"}



def init_glpi_session(username, password):
    url = f"{API_URL}/initSession"
    headers = {
        "App-Token": APP_TOKEN,
        "Content-Type": "application/json"
    }
    payload = {"login": username, "password": password}
    try:
        response = requests.post(url, json=payload, headers=headers)
        if response.status_code == 200:
            return response.json().get("session_token")
    except Exception as e:
        print("[SESSION ERROR]", e)
    return None

def get_glpi_entities(session_token):
    url = f"{API_URL}/Entity"
    headers = {
        "App-Token": APP_TOKEN,
        "Session-Token": session_token
    }
    try:
        response = requests.get(url, headers=headers)
        return response.json() if response.status_code == 200 else []
    except Exception as e:
        print("[ENTITIES ERROR]", e)
        return []

def get_glpi_tickets_by_entity(entity_id, session_token):
    url = (
        f"{API_URL}/Ticket?"
        f"criteria[0][field]=entities_id&criteria[0][searchtype]=equals&criteria[0][value]={entity_id}&"
        f"is_recursive=0"
    )

    headers = {
        "App-Token": APP_TOKEN,
        "Session-Token": session_token
    }

    try:
        response = requests.get(url, headers=headers)
        tickets = response.json()

        # filtrage manuel
        tickets = [t for t in tickets if str(t.get("entities_id")) == str(entity_id)]

        detailed_tickets = []
        for ticket in tickets:
            ticket_id = ticket.get("id")
            if not ticket_id:
                continue
            # Requête pour avoir les détails complets du ticket
            detail_url = f"{API_URL}/Ticket/{ticket_id}"
            detail_resp = requests.get(detail_url, headers=headers)
            if detail_resp.status_code == 200:
                detailed_ticket = detail_resp.json()
                detailed_tickets.append(detailed_ticket)

        return {"tickets": detailed_tickets}
    except Exception as e:
        return {"error": f"Erreur API : {e}"}


def summarize_tickets(tickets):
    status_count = Counter()
    priority_count = Counter()
    last_modification = None
    ttrs = []

    for ticket in tickets:
        status_count[ticket.get('status')] += 1
        priority_count[ticket.get('priority')] += 1

        mod_date = ticket.get('date_mod')
        if mod_date:
            mod_dt = datetime.strptime(mod_date, "%Y-%m-%d %H:%M:%S")
            if last_modification is None or mod_dt > last_modification:
                last_modification = mod_dt

        if ticket.get('solvedate') and ticket.get('date'):
            try:
                solved = datetime.strptime(ticket['solvedate'], "%Y-%m-%d %H:%M:%S")
                opened = datetime.strptime(ticket['date'], "%Y-%m-%d %H:%M:%S")
                ttrs.append((solved - opened).total_seconds())
            except:
                continue

    avg_ttr = sum(ttrs) / len(ttrs) / 60 if ttrs else 0

    return {
        "nb_total": len(tickets),
        "statuts": dict(status_count),
        "priorites": dict(priority_count),
        "last_mod": last_modification.strftime("%Y-%m-%d %H:%M:%S") if last_modification else "N/A",
        "ttr_moyen": f"{avg_ttr:.1f} minutes"
    }
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx import Presentation

# Fonction pour ajouter un tableau de tickets à une slide
def add_ticket_table(slide, tickets):
    rows = len(tickets) + 1
    cols = 5

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 + 0.3 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # En-têtes
    table.cell(0, 0).text = "ID"
    table.cell(0, 1).text = "Nom"
    table.cell(0, 2).text = "Statut"
    table.cell(0, 3).text = "Urgence"
    table.cell(0, 4).text = "Création"

    # Données
    for i, ticket in enumerate(tickets):
        table.cell(i + 1, 0).text = str(ticket.get('id', 'N/A'))
        table.cell(i + 1, 1).text = ticket.get('name', 'Sans nom')
        table.cell(i + 1, 2).text = get_ticket_status(ticket.get('status'))
        table.cell(i + 1, 3).text = get_ticket_priority(ticket.get('priority'))
        table.cell(i + 1, 4).text = format_ticket_date(ticket.get('date_creation'))

    # Style du texte
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = RGBColor(50, 50, 50)

    # Couleur de fond verte claire
    color_table_headers(table, RGBColor(95, 200, 192))
    apply_alternate_row_colors(table)  # <-- ici







from collections import defaultdict
from datetime import datetime

def group_tickets_by_month(tickets):
    grouped = defaultdict(list)
    for ticket in tickets:
        date_str = ticket.get("date")
        if date_str:
            try:
                dt = datetime.strptime(date_str, "%Y-%m-%d %H:%M:%S")
                key = dt.strftime("%Y-%m")  # format: "2025-04"
                grouped[key].append(ticket)
            except:
                continue
    return grouped
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ✅ Fonction : colorer seulement les en-têtes du tableau
def color_table_headers(table, rgb_color):
    header_row = table.rows[0]
    for cell in header_row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb_color
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = RGBColor(0, 0, 0)

from pptx.dml.color import RGBColor

def apply_alternate_row_colors(table, color1=RGBColor(234, 244, 243), color2=RGBColor(245, 250, 249)):
    """
    Applique un fond de couleur alternée (ligne paire = color1, ligne impaire = color2)
    aux lignes du tableau, en ignorant la première ligne (en-têtes).
    """
    rows = len(table.rows)
    cols = len(table.columns)

    for i in range(1, rows):  # Commence à 1 pour sauter les en-têtes
        for j in range(cols):
            cell = table.cell(i, j)
            cell.fill.solid()
            if i % 2 == 0:
                cell.fill.fore_color.rgb = color1
            else:
                cell.fill.fore_color.rgb = color2





def get_ticket_status(code):
    return {
        1: "Nouveau",
        2: "En cours",
        3: "En attente",
        4: "Résolu",
        5: "Clos"
    }.get(code, "Inconnu")

def get_ticket_priority(code):
    return {
        1: "Très basse",
        2: "Basse",
        3: "Moyenne",
        4: "Haute",
        5: "Très haute"
    }.get(code, "Inconnue")

def format_ticket_date(date_string):
    try:
        date_obj = datetime.strptime(date_string, '%Y-%m-%d %H:%M:%S')
        return date_obj.strftime('%d/%m/%Y %H:%M')
    except:
        return "Date invalide"

def get_users_by_entity(entity_id, session_token):
    url = f"{API_URL}/User"
    headers = {
        "App-Token": APP_TOKEN,
        "Session-Token": session_token
    }
    params = {
        "criteria[0][field]": "entities_id",
        "criteria[0][searchtype]": "equals",
        "criteria[0][value]": entity_id
    }

    try:
        response = requests.get(url, headers=headers, params=params)
        if response.status_code == 200:
            users = response.json()

            # Filtrage manuel pour vérifier si l'ID de l'entité correspond à celle demandée
            filtered_users = [user for user in users if str(user.get('entities_id')) == str(entity_id)]

            return [user.get("name", "Utilisateur inconnu") for user in filtered_users]
        else:
            return []
    except Exception as e:
        print(f"[ERROR] Failed to retrieve users: {e}")
        return []
